"""
Generates demo/slides.pptx with fully editable text using python-pptx.
Run from repo root: python demo/build_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0d, 0x11, 0x17)   # slide background
BLUE      = RGBColor(0x58, 0xa6, 0xff)   # headings / accent
WHITE     = RGBColor(0xf0, 0xf6, 0xfc)   # body text
GREY      = RGBColor(0x8b, 0x94, 0x9e)   # subdued text
GREEN     = RGBColor(0x3f, 0xb9, 0x50)   # positive / code
ORANGE    = RGBColor(0xf0, 0x88, 0x3e)   # warning
RED       = RGBColor(0xff, 0x7b, 0x72)   # error / problem
PURPLE    = RGBColor(0xbc, 0x8c, 0xff)   # copilot
CARD_BG   = RGBColor(0x16, 0x1b, 0x22)   # card background
BORDER    = RGBColor(0x30, 0x36, 0x3d)   # borders

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def add_box(slide, x, y, w, h, fill=CARD_BG, line=BORDER):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Emu(9525)   # 0.75pt
    return shape

def add_box_with_text(slide, text, x, y, w, h, size=12, bold=False,
                      color=WHITE, fill=CARD_BG, line=BORDER,
                      align=PP_ALIGN.LEFT):
    add_box(slide, x, y, w, h, fill=fill, line=line)
    pad = Inches(0.12)
    add_text(slide, text, x+pad, y+pad, w-pad*2, h-pad*2,
             size=size, bold=bold, color=color, align=align)

def label(slide, text, x, y, w=Inches(10)):
    add_text(slide, text, x, y, w, Inches(0.3),
             size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

def hline(slide, y):
    line = slide.shapes.add_connector(1, Inches(0.5), y, Inches(12.83), y)
    line.line.color.rgb = BORDER
    line.line.width = Emu(9525)

def conn(slide, x1, y1, x2, y2, color=BLUE, lw=Emu(15876), ctype=1):
    """Draw a connector (ctype 1=straight, 2=elbow/angular) with arrowhead at (x2,y2)."""
    from lxml import etree
    c = slide.shapes.add_connector(ctype, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = lw
    # Arrowhead via XML on the connector's spPr/ln element
    try:
        from pptx.oxml.ns import qn
        sp_pr = c._element.find(qn('p:spPr'))
        ans   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ln    = sp_pr.find(f'{{{ans}}}ln')
        if ln is None:
            ln = etree.SubElement(sp_pr, f'{{{ans}}}ln')
        tail = etree.SubElement(ln, f'{{{ans}}}tailEnd')
        tail.set('type', 'arrow')
        tail.set('w', 'med')
        tail.set('len', 'med')
    except Exception:
        pass   # arrowhead failed gracefully — line still visible
    return c

def conn_label(slide, x1, y1, x2, y2, text, color=BLUE):
    """Small label centred on a connector."""
    mx, my = (x1 + x2) / 2, (y1 + y2) / 2
    add_text(slide, text, Inches(mx - 0.9), Inches(my - 0.17),
             Inches(1.8), Inches(0.32), size=8, color=color, align=PP_ALIGN.CENTER)

def seq_badge(slide, n, x, y, color=BLUE):
    """Small filled square with a sequence number, placed at (x, y)."""
    add_box(slide, Inches(x - 0.14), Inches(y - 0.14), Inches(0.28), Inches(0.28),
            fill=color, line=color)
    add_text(slide, str(n), Inches(x - 0.14), Inches(y - 0.16), Inches(0.28), Inches(0.28),
             size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

# ── Slide builder functions ───────────────────────────────────────────────────

def slide_title(prs):
    s = new_slide(prs)
    label(s, "AGENTIC AI", Inches(0.5), Inches(1.5))
    add_text(s, "Multi-Repo Agentic\nOrchestration Framework",
             Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.8),
             size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "Cross-repo AI awareness for lithography software development",
             Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5),
             size=18, color=GREY, align=PP_ALIGN.CENTER)

    tech = ["MCP Protocol", "RAG", "ChromaDB", "sentence-transformers", "FastMCP", "GitHub Copilot", "Python 3.11+"]
    x = Inches(1.2)
    for t in tech:
        bw = Inches(1.5)
        add_box_with_text(s, t, x, Inches(4.8), bw, Inches(0.35),
                          size=10, bold=True, color=BLUE,
                          fill=RGBColor(0x1f,0x3a,0x5f), line=RGBColor(0x2d,0x59,0x86),
                          align=PP_ALIGN.CENTER)
        x += bw + Inches(0.08)


def slide_problem(prs):
    s = new_slide(prs)
    label(s, "THE PROBLEM", Inches(0.5), Inches(0.3))
    add_text(s, "Copilot sees only one repo at a time",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Real features span multiple lithography subsystems — but your AI assistant is blind to them",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
             size=14, color=GREY, align=PP_ALIGN.CENTER)

    repos = ["scan_manager", "illumination", "expose_sequence", "wafer_handler"]
    positions = [Inches(0.4), Inches(3.4), Inches(6.4), Inches(9.4)]
    for name, x in zip(repos, positions):
        add_box_with_text(s, f"📦  {name}", x, Inches(1.9), Inches(2.8), Inches(0.7),
                          size=13, bold=True, color=WHITE)

    add_box_with_text(s, "🤖  GitHub Copilot\n     sees only scan_manager",
                      Inches(3.4), Inches(2.9), Inches(2.8), Inches(0.8),
                      size=12, color=PURPLE,
                      fill=RGBColor(0x2d,0x1f,0x47), line=PURPLE)

    problems = [
        "🔴   A dose control change in expose_sequence must align with scan_manager timing — Copilot cannot see this",
        "🔴   Illumination parameter updates affect exposure logic across repos — completely invisible to the AI",
        "🔴   Developer must manually gather context from all subsystem repos before every cross-repo question",
        "🔴   AI suggestions are incomplete — only one subsystem's perspective, missing critical dependencies",
    ]
    y = Inches(4.0)
    for p in problems:
        add_box_with_text(s, p, Inches(0.4), y, Inches(12.5), Inches(0.5),
                          size=12, color=WHITE,
                          fill=RGBColor(0x1a,0x0a,0x0a), line=RED)
        y += Inches(0.6)


def slide_solution(prs):
    s = new_slide(prs)
    label(s, "THE SOLUTION", Inches(0.5), Inches(0.3))
    add_text(s, "An orchestra of specialised agents",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Each subsystem repo has its own agent. An orchestrator routes intelligently. Copilot synthesises.",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
             size=14, color=GREY, align=PP_ALIGN.CENTER)

    # Copilot box at top
    add_box_with_text(s, "🤖  GitHub Copilot  (Agent mode)",
                      Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.6),
                      size=14, bold=True, color=PURPLE,
                      fill=RGBColor(0x2d,0x1f,0x47), line=PURPLE, align=PP_ALIGN.CENTER)

    add_text(s, "↕", Inches(6.4), Inches(2.5), Inches(0.5), Inches(0.4),
             size=20, color=BLUE, align=PP_ALIGN.CENTER)

    # Orchestrator
    add_box_with_text(s, "🎯  Orchestrator\n     Routes by RAG relevance score\n     Returns top 2 relevant repos",
                      Inches(4.2), Inches(3.0), Inches(4.9), Inches(0.95),
                      size=12, bold=False, color=GREEN,
                      fill=RGBColor(0x0d,0x1f,0x12), line=GREEN, align=PP_ALIGN.CENTER)

    add_text(s, "↙                    ↓                    ↘",
             Inches(2.5), Inches(4.05), Inches(8.0), Inches(0.4),
             size=16, color=BLUE, align=PP_ALIGN.CENTER)

    # Repo agents
    agents = ["📦  scan_manager\n     RAG index", "📦  illumination\n     RAG index", "📦  expose_sequence\n     RAG index"]
    positions = [Inches(0.8), Inches(4.9), Inches(9.0)]
    for name, x in zip(agents, positions):
        add_box_with_text(s, name, x, Inches(4.55), Inches(3.3), Inches(0.75),
                          size=12, color=BLUE,
                          fill=RGBColor(0x0d,0x1f,0x35), line=BLUE, align=PP_ALIGN.CENTER)

    add_box_with_text(s, "✅  Feature Analysis Document — cross-subsystem, complete, grounded in real knowledge",
                      Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
                      size=13, bold=True, color=GREEN,
                      fill=RGBColor(0x0d,0x1f,0x12), line=GREEN, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    s = new_slide(prs)
    label(s, "ARCHITECTURE", Inches(0.5), Inches(0.3))
    add_text(s, "How it all fits together",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    headers = ["Layer", "Component", "Role"]
    col_x   = [Inches(0.4), Inches(2.8), Inches(7.3)]
    col_w   = [Inches(2.2), Inches(4.3), Inches(5.6)]

    rows = [
        ("AI Interface",   "GitHub Copilot + Agent mode",          "Developer's entry point — types feature request"),
        ("Instructions",   "copilot-instructions.md",              "Tells Copilot to call MCP tools in order"),
        ("Tool Config",    ".vscode/mcp.json  (auto-generated)",   "Wires up all MCP servers with absolute paths"),
        ("Orchestrator",   "router_mcp_server.py",                 "Routes query to most relevant subsystem repos"),
        ("Repo Agents",    "mcp_server.py  (per repo)",            "Exposes query_repo tool, runs RAG search"),
        ("Knowledge",      "repo_rag.py + ChromaDB",               "Indexes docs + source code, serves queries"),
        ("Embedding",      "all-MiniLM-L6-v2  (local, ~90 MB)",   "Converts text to semantic vectors, fully offline"),
    ]

    # Header row
    y = Inches(1.4)
    for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_box_with_text(s, hdr, cx, y, cw, Inches(0.38),
                          size=12, bold=True, color=BLUE,
                          fill=RGBColor(0x16,0x1b,0x22), line=BORDER, align=PP_ALIGN.CENTER)

    # Data rows
    y += Inches(0.38)
    for row in rows:
        for i, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
            c = BLUE if i == 0 else (GREEN if i == 1 else GREY)
            add_box_with_text(s, val, cx, y, cw, Inches(0.55),
                              size=11, color=c,
                              fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.55)

    add_text(s, "★  One setup.py generates all mcp.json and copilot-instructions.md files automatically",
             Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
             size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_tech_stack(prs):
    s = new_slide(prs)
    label(s, "TECHNOLOGY STACK", Inches(0.5), Inches(0.3))
    add_text(s, "Purpose-built from proven components",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Every layer is local — no cloud calls, no data leaves your machine",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.35),
             size=14, color=GREY, align=PP_ALIGN.CENTER)

    layers = [
        (PURPLE, "AI Interface",     "GitHub Copilot  ·  Agent Mode  ·  copilot-instructions.md"),
        (BLUE,   "Communication",    "MCP — Model Context Protocol  ·  stdio transport  ·  JSON-RPC 2.0  ·  FastMCP"),
        (GREEN,  "Knowledge / RAG",  "ChromaDB (vector store)  ·  sentence-transformers  ·  all-MiniLM-L6-v2 (local model)"),
        (ORANGE, "Config / Infra",   "config.yaml  ·  mcp.json (auto-generated)  ·  Python 3.11+  ·  pysqlite3-binary"),
        (RED,    "IDE",              "VS Code  ·  .vscode/mcp.json  ·  .github/copilot-instructions.md"),
    ]

    y = Inches(1.75)
    for color, layer_name, items in layers:
        add_box_with_text(s, layer_name, Inches(0.4), y, Inches(2.2), Inches(0.55),
                          size=12, bold=True, color=color,
                          fill=RGBColor(0x16,0x1b,0x22), line=color, align=PP_ALIGN.CENTER)
        add_box_with_text(s, items, Inches(2.7), y, Inches(10.2), Inches(0.55),
                          size=12, color=WHITE,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.65)


def slide_workflow(prs):
    s = new_slide(prs)
    label(s, "WORKFLOW", Inches(0.5), Inches(0.3))
    add_text(s, "4 steps, fully automatic",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 'Developer: "Add real-time dose control feedback during wafer exposure"',
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
             size=14, color=ORANGE, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Route",          "orchestrator_router\n.get_relevant_repos()",
               "Queries each repo RAG\nScores by content length\nReturns top 2 repos",
               "expose_sequence: 0.91\nillumination: 0.74"),
        ("2", "Own Repo",       "scan_manager\n.query_repo()",
               "Searches own docs\n+ source code via RAG",
               "Returns scan timing,\nstage control knowledge"),
        ("3", "Query Peers",    "expose_sequence\n.query_repo()\nillumination\n.query_repo()",
               "Retrieves dose calc,\npulse control knowledge\nfrom selected peers",
               "Returns combined\nknowledge chunks"),
        ("4", "Synthesise",     "Feature Analysis\nDocument",
               "Current State per repo\n+ cross-subsystem\nSolution Design",
               "Complete cross-repo\nanswer for Copilot"),
    ]

    x = Inches(0.3)
    for num, title, tool, body, result in steps:
        # Step number circle (simulated with box)
        add_box_with_text(s, num, x+Inches(0.85), Inches(1.85), Inches(0.45), Inches(0.45),
                          size=14, bold=True, color=BLUE,
                          fill=RGBColor(0x1f,0x3a,0x5f), line=BLUE, align=PP_ALIGN.CENTER)
        # Main card
        add_box_with_text(s, title, x, Inches(2.4), Inches(2.9), Inches(0.45),
                          size=13, bold=True, color=WHITE,
                          fill=RGBColor(0x16,0x1b,0x22), line=BLUE, align=PP_ALIGN.CENTER)
        add_box_with_text(s, tool, x, Inches(2.95), Inches(2.9), Inches(0.7),
                          size=10, bold=True, color=GREEN,
                          fill=RGBColor(0x0d,0x1f,0x12), line=GREEN, align=PP_ALIGN.CENTER)
        add_box_with_text(s, body, x, Inches(3.75), Inches(2.9), Inches(1.1),
                          size=11, color=GREY,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        add_box_with_text(s, result, x, Inches(4.95), Inches(2.9), Inches(0.7),
                          size=11, color=ORANGE,
                          fill=RGBColor(0x1f,0x12,0x00), line=ORANGE)
        if num != "4":
            add_text(s, "→", x+Inches(2.93), Inches(3.5), Inches(0.35), Inches(0.4),
                     size=20, bold=True, color=BLUE)
        x += Inches(3.1) + Inches(0.17)


def slide_rag(prs):
    s = new_slide(prs)
    label(s, "UNDER THE HOOD — RAG", Inches(0.5), Inches(0.3))
    add_text(s, "How each agent knows its subsystem",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Retrieval-Augmented Generation — semantic search over your own docs and source code",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.35),
             size=14, color=GREY, align=PP_ALIGN.CENTER)

    # Phase 1
    add_box_with_text(s, "PHASE 1 — Indexing  (runs once on first start)",
                      Inches(0.4), Inches(1.75), Inches(5.9), Inches(0.4),
                      size=12, bold=True, color=BLUE,
                      fill=RGBColor(0x1f,0x3a,0x5f), line=BLUE)
    phase1 = [
        "📄  .md docs (knowledge/ + BB-*/*/docs)  →  split into chunks  →  embed  →  ChromaDB",
        "💻  Source code (BB-*/*/com)              →  split into chunks  →  embed  →  ChromaDB",
        "💾  Stored in .chroma_db/  —  persistent cache, rebuilt on each test run",
    ]
    y = Inches(2.25)
    for line in phase1:
        add_box_with_text(s, line, Inches(0.4), y, Inches(12.5), Inches(0.42),
                          size=11, color=WHITE,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.47)

    # Phase 2
    add_box_with_text(s, "PHASE 2 — Query  (every call to query_repo)",
                      Inches(0.4), Inches(3.7), Inches(5.9), Inches(0.4),
                      size=12, bold=True, color=GREEN,
                      fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)
    phase2 = [
        '❓  Feature request: "dose control feedback wafer exposure scan timing"',
        "     →  embed query  →  search docs collection  (top 3 semantically similar chunks)",
        "     →  search code collection  (top 3 chunks) — always, not just when docs are thin",
        "     →  return:  [From documentation] ...  +  [From source code] ...",
        "     →  orchestrator scores:  min(content_length / 800,  1.0)",
    ]
    y = Inches(4.2)
    for line in phase2:
        add_box_with_text(s, line, Inches(0.4), y, Inches(12.5), Inches(0.42),
                          size=11, color=WHITE if not line.startswith("❓") else ORANGE,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.47)


def slide_mcp(prs):
    s = new_slide(prs)
    label(s, "UNDER THE HOOD — MCP", Inches(0.5), Inches(0.3))
    add_text(s, "How agents communicate",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Model Context Protocol — stdin/stdout only, no network ports, no cloud",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.35),
             size=14, color=GREY, align=PP_ALIGN.CENTER)

    protocol = [
        ('msg 1', 'initialize',                    'Handshake — agree on protocol version'),
        ('msg 2', 'notifications/initialized',     'Client confirms ready'),
        ('msg 3', 'tools/call  →  query_repo',     'Actual RAG query — returns relevant knowledge'),
        ('reply', 'result.content[0].text',        'RELEVANT KNOWLEDGE: [docs] ... [source code] ...'),
    ]
    y = Inches(1.75)
    for tag, method, desc in protocol:
        add_box_with_text(s, tag, Inches(0.4), y, Inches(0.9), Inches(0.45),
                          size=10, bold=True, color=GREY,
                          fill=RGBColor(0x16,0x1b,0x22), line=BORDER, align=PP_ALIGN.CENTER)
        add_box_with_text(s, method, Inches(1.4), y, Inches(3.8), Inches(0.45),
                          size=11, bold=True, color=GREEN,
                          fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)
        add_box_with_text(s, desc, Inches(5.3), y, Inches(7.6), Inches(0.45),
                          size=11, color=GREY,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.52)

    # Two server cards
    add_box_with_text(s,
        "router_mcp_server.py\n\nTool: get_relevant_repos\n\nCalled by Copilot in step 1\nSpawns each repo agent as subprocess\nScores responses, returns top 2",
        Inches(0.4), Inches(4.3), Inches(6.1), Inches(2.8),
        size=12, color=WHITE, fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)

    add_box_with_text(s,
        "mcp_server.py  (per subsystem repo)\n\nTool: query_repo\n\nCalled by Copilot + orchestrator\nSearches ChromaDB for relevant chunks\nReturns docs + source code knowledge",
        Inches(6.8), Inches(4.3), Inches(6.1), Inches(2.8),
        size=12, color=WHITE, fill=RGBColor(0x0d,0x1f,0x35), line=BLUE)


def slide_setup(prs):
    s = new_slide(prs)
    label(s, "SETUP", Inches(0.5), Inches(0.3))
    add_text(s, "What the user does — once",
             Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Clone framework",       "git clone <this repo>"),
        ("2", "Copy embedding model",  "Copy models/all-MiniLM-L6-v2/ into repo root  (offline, ~90 MB)"),
        ("3", "Install dependencies",  "pip install -r requirements.txt"),
        ("4", "Configure orchestrator","Edit orchestrator/mcp/config.yaml — list all subsystem repos + paths"),
        ("5", "Add agent files",       "cp repo-agent/mcp/ knowledge/ .github/  into each subsystem repo"),
        ("6", "Configure each repo",   "Edit each repo's mcp/config.yaml — set repo_name and src_path (glob ok)"),
        ("7", "Run setup",             "python orchestrator/setup.py — generates mcp.json + copilot-instructions.md"),
        ("8", "Open in VS Code",       "source .venv/bin/activate  &&  code /path/to/repo"),
    ]

    y = Inches(1.4)
    for num, title, detail in steps:
        add_box_with_text(s, num, Inches(0.4), y, Inches(0.38), Inches(0.48),
                          size=12, bold=True, color=BLUE,
                          fill=RGBColor(0x1f,0x3a,0x5f), line=BLUE, align=PP_ALIGN.CENTER)
        add_box_with_text(s, title, Inches(0.88), y, Inches(2.6), Inches(0.48),
                          size=12, bold=True, color=WHITE,
                          fill=RGBColor(0x16,0x1b,0x22), line=BORDER)
        add_box_with_text(s, detail, Inches(3.58), y, Inches(9.7), Inches(0.48),
                          size=11, color=GREY,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.55)

    add_text(s, "★  Only config.yaml files are ever edited by the user — everything else is generated",
             Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
             size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_demo(prs):
    s = new_slide(prs)
    label(s, "LIVE DEMO", Inches(0.5), Inches(0.5))
    add_text(s, "Let's see it in action",
             Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.6),
             size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    cmds = [
        ("Terminal 1 — watch live logs",
         "tail -f /tmp/mcp-orchestration.log"),
        ("Terminal 2 — run the simulation",
         'python test_mcp.py "dose control feedback wafer exposure scan timing"'),
    ]
    y = Inches(2.0)
    for title, cmd in cmds:
        add_text(s, title, Inches(1.0), y, Inches(11.3), Inches(0.35),
                 size=13, bold=True, color=GREY)
        add_box_with_text(s, cmd, Inches(1.0), y+Inches(0.38), Inches(11.3), Inches(0.5),
                          size=13, bold=True, color=GREEN,
                          fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)
        y += Inches(1.1)

    add_text(s, "Output", Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.35),
             size=13, bold=True, color=GREY)
    outputs = [
        "LOG  |  [scan_manager]  [RAG]   Searching docs (top_k=3)...",
        "LOG  |  [scan_manager]  [RAG]   Docs: 3 results, 480 chars",
        "LOG  |  [scan_manager]  [RAG]   Searching source code (top_k=3)...",
        "LOG  |  [scan_manager]  [RAG]   Code: 3 results, 620 chars",
        "LOG  |  [orchestrator]  [RESULT] Selected: ['expose_sequence', 'illumination']",
    ]
    y = Inches(4.7)
    for line in outputs:
        add_box_with_text(s, line, Inches(1.0), y, Inches(11.3), Inches(0.37),
                          size=10, color=GREEN,
                          fill=RGBColor(0x0d,0x11,0x17), line=BORDER)
        y += Inches(0.4)

    add_box_with_text(s, "🎯  Output saved to feature_analysis.md — paste into Copilot chat for cross-subsystem Solution Design",
                      Inches(1.0), Inches(7.05), Inches(11.3), Inches(0.38),
                      size=13, bold=True, color=GREEN,
                      fill=RGBColor(0x0d,0x1f,0x12), line=GREEN, align=PP_ALIGN.CENTER)


def slide_rag_init_flow(prs):
    """Phase 1 — RAG index creation per MCP server (cold start)."""
    s = new_slide(prs)

    # ── Header ─────────────────────────────────────────────────────────────────
    add_text(s, "DATA FLOW — PHASE 1  ·  INDEX CREATION",
             Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.26),
             size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "RAG Index Built Once per MCP Server on Cold Start",
             Inches(0.4), Inches(0.35), Inches(12.5), Inches(0.42),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hline(s, Inches(0.83))

    # ── Box geometry (all coords in inches) ──────────────────────────────────────
    # Slide center x = 13.33/2 = 6.665
    TX, TY, TW, TH = 0.20, 1.00, 2.60, 0.75   # test_mcp.py
    MX, MY, MW, MH = 4.20, 1.00, 5.30, 0.75   # mcp_server.py
    DX, DY, DW, DH = 1.50, 2.38, 3.20, 0.78   # knowledge/*.md
    SX, SY, SW, SH = 8.45, 2.38, 3.45, 0.78   # ./src/ source code
    KX, KY, KW, KH = 3.00, 3.58, 7.30, 1.00   # Chunker
    EX, EY, EW, EH = 4.05, 5.00, 5.25, 0.85   # Embedding Model
    VX, VY, VW, VH = 3.50, 6.18, 6.30, 0.78   # ChromaDB  (V for Vector store)

    # derived midpoints / edges
    T_r   = TX + TW;   T_cy  = TY + TH/2
    M_l   = MX;        M_cx  = MX + MW/2;  M_bot = MY + MH
    D_cx  = DX + DW/2; D_bot = DY + DH
    S_cx  = SX + SW/2; S_bot = SY + SH
    K_cx  = KX + KW/2; K_bot = KY + KH
    E_cx  = EX + EW/2; E_top = EY;         E_bot = EY + EH
    V_cx  = VX + VW/2; V_top = VY

    # ── Component boxes ────────────────────────────────────────────────────────

    def bx(slide, title, sub, x, y, w, h, fill, ln, tc=WHITE, sc=None):
        add_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=ln)
        p = 0.10
        add_text(slide, title, Inches(x+p), Inches(y+p), Inches(w-2*p), Inches(h*0.48),
                 size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
        if sub:
            add_text(slide, sub, Inches(x+p), Inches(y+h*0.52), Inches(w-2*p), Inches(h*0.42),
                     size=8, color=sc or GREY, align=PP_ALIGN.CENTER)

    bx(s, "test_mcp.py", "MCP client  ·  subprocess.run()",
       TX, TY, TW, TH, RGBColor(0x0d,0x1f,0x0d), GREEN, GREEN)

    bx(s, "mcp_server.py  ·  RepoRAG.build_or_load_index()",
       "checks collection exists → builds if absent → loads if cached",
       MX, MY, MW, MH, RGBColor(0x0d,0x1f,0x35), BLUE, BLUE)

    bx(s, "knowledge/*.md", "documentation · markdown files",
       DX, DY, DW, DH, RGBColor(0x16,0x1b,0x22), GREY, WHITE)

    bx(s, "./src/  source code", ".py  .cpp  .h  .java  .go  .rs  ···  12 languages",
       SX, SY, SW, SH, RGBColor(0x16,0x1b,0x22), GREY, WHITE)

    add_box(s, Inches(KX), Inches(KY), Inches(KW), Inches(KH),
            fill=RGBColor(0x22,0x14,0x00), line=ORANGE)
    add_text(s, "Chunker  (repo_rag.py)",
             Inches(KX+0.10), Inches(KY+0.10), Inches(KW-0.20), Inches(0.30),
             size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, "Docs → paragraph split     ·     Code → function / class boundary split\n"
                "max_chunk_chars = 1500   ·   chunk_overlap_chars = 150   ·   min chunk = 40 chars",
             Inches(KX+0.14), Inches(KY+0.44), Inches(KW-0.28), Inches(0.50),
             size=10, color=GREY, align=PP_ALIGN.CENTER)

    add_box(s, Inches(EX), Inches(EY), Inches(EW), Inches(EH),
            fill=RGBColor(0x10,0x18,0x25), line=BLUE)
    add_text(s, "Embedding Model  ·  all-MiniLM-L6-v2  ·  local  ·  offline",
             Inches(EX+0.10), Inches(EY+0.10), Inches(EW-0.20), Inches(0.30),
             size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "encode(texts, normalize_embeddings=True, batch_size=512)  ·  22.7 M params  ·  ~300 MB RAM\n"
                "text chunks  →  384-dim float32 vectors  (encoder-only: one-way, no inverse)",
             Inches(EX+0.14), Inches(EY+0.44), Inches(EW-0.28), Inches(0.36),
             size=10, color=GREY, align=PP_ALIGN.CENTER)

    add_box(s, Inches(VX), Inches(VY), Inches(VW), Inches(VH),
            fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)
    add_text(s, "ChromaDB  ·  .chroma_db/  ·  SQLite WAL",
             Inches(VX+0.10), Inches(VY+0.10), Inches(VW-0.20), Inches(0.28),
             size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, "collection.add(documents=text_chunks, embeddings=vectors, ids=[…])\n"
                "docs_collection  +  code_collection  ·  persistent cache  ·  concurrent reads safe",
             Inches(VX+0.14), Inches(VY+0.38), Inches(VW-0.28), Inches(0.34),
             size=10, color=GREY, align=PP_ALIGN.CENTER)

    # ── Angular connectors + sequence badges + data labels ────────────────────

    def alabel(slide, text, x, y, col=GREY):
        add_text(slide, text, Inches(x), Inches(y), Inches(2.5), Inches(0.20),
                 size=7, color=col, align=PP_ALIGN.CENTER)

    # ① test_mcp → mcp_server  (horizontal)
    conn(s, T_r, T_cy, M_l, T_cy, GREEN, ctype=2)
    seq_badge(s, 1, (T_r+M_l)/2, T_cy-0.28, GREEN)
    alabel(s, "subprocess.run(mcp_server.py)  ·  JSON-RPC initialize",
           (T_r+M_l)/2 - 1.25, T_cy - 0.46, GREEN)

    # ② mcp_server → docs (elbow down-left)
    conn(s, M_cx-1.5, M_bot, D_cx, DY, GREY, ctype=2)
    seq_badge(s, 2, (M_cx-1.5+D_cx)/2 - 0.1, (M_bot+DY)/2, GREY)
    alabel(s, "reads .md files", (M_cx-1.5+D_cx)/2 - 1.25, (M_bot+DY)/2 - 0.22, GREY)

    # ② mcp_server → src  (elbow down-right)
    conn(s, M_cx+1.5, M_bot, S_cx, SY, GREY, ctype=2)
    alabel(s, "reads source files", (M_cx+1.5+S_cx)/2 - 1.25, (M_bot+SY)/2 - 0.22, GREY)

    # ③ docs → chunker left  (elbow down-right)
    conn(s, D_cx, D_bot, KX+0.6, KY, ORANGE, ctype=2)
    seq_badge(s, 3, (D_cx+KX+0.6)/2 - 0.1, (D_bot+KY)/2, ORANGE)
    alabel(s, "paragraph chunks", (D_cx+KX+0.6)/2 - 1.25, (D_bot+KY)/2 - 0.22, ORANGE)

    # ③ src → chunker right  (elbow down-left)
    conn(s, S_cx, S_bot, KX+KW-0.6, KY, ORANGE, ctype=2)
    alabel(s, "fn-boundary chunks", (S_cx+KX+KW-0.6)/2 - 1.25, (S_bot+KY)/2 - 0.22, ORANGE)

    # ④ chunker → embedding  (straight down)
    conn(s, K_cx, K_bot, E_cx, E_top, BLUE, ctype=1)
    seq_badge(s, 4, K_cx+0.22, (K_bot+E_top)/2, BLUE)
    alabel(s, "text chunks  (batch 512)", K_cx+0.38, (K_bot+E_top)/2 - 0.12, BLUE)

    # ⑤ embedding → chromadb  (straight down)
    conn(s, E_cx, E_bot, V_cx, V_top, GREEN, ctype=1)
    seq_badge(s, 5, E_cx+0.22, (E_bot+V_top)/2, GREEN)
    alabel(s, "384-dim float32[]  +  original text  +  ids", E_cx+0.38, (E_bot+V_top)/2 - 0.12, GREEN)


def slide_rag_query_flow(prs):
    """Phase 2 — full test_mcp.py query flow: routing + per-repo RAG search."""
    s = new_slide(prs)

    add_text(s, "DATA FLOW — PHASE 2  ·  QUERY  (test_mcp.py)",
             Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.26),
             size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "Feature Request → Routing → Per-Repo RAG Search → Feature Analysis Document",
             Inches(0.4), Inches(0.35), Inches(12.5), Inches(0.42),
             size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hline(s, Inches(0.83))

    # ── Phase lane backgrounds ─────────────────────────────────────────────────
    add_box(s, Inches(0.15), Inches(0.90), Inches(13.05), Inches(2.52),
            fill=RGBColor(0x0d,0x14,0x0d), line=RGBColor(0x1a,0x30,0x1a))
    add_text(s, "STEP 1 — ROUTING  (router_mcp_server.py queries each peer repo to score relevance)",
             Inches(0.28), Inches(0.93), Inches(12.8), Inches(0.24),
             size=9, bold=True, color=GREEN)

    add_box(s, Inches(0.15), Inches(3.52), Inches(13.05), Inches(3.28),
            fill=RGBColor(0x0d,0x14,0x1f), line=RGBColor(0x1a,0x28,0x44))
    add_text(s, "STEP 2 — QUERY  (per-repo RAG: embed query → ChromaDB similarity search → top-k chunks)",
             Inches(0.28), Inches(3.55), Inches(12.8), Inches(0.24),
             size=9, bold=True, color=BLUE)

    # ── ROUTING boxes ──────────────────────────────────────────────────────────
    # test_mcp.py (routing)
    add_box(s, Inches(0.28), Inches(1.22), Inches(2.10), Inches(0.85),
            fill=RGBColor(0x0d,0x1f,0x0d), line=GREEN)
    add_text(s, "test_mcp.py", Inches(0.38), Inches(1.30), Inches(1.90), Inches(0.36),
             size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, "MCP client", Inches(0.38), Inches(1.66), Inches(1.90), Inches(0.22),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # router_mcp_server.py
    add_box(s, Inches(3.20), Inches(1.22), Inches(4.30), Inches(0.85),
            fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)
    add_text(s, "router_mcp_server.py", Inches(3.30), Inches(1.30), Inches(4.10), Inches(0.30),
             size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, "get_relevant_repos(requesting_repo, feature_description)",
             Inches(3.30), Inches(1.62), Inches(4.10), Inches(0.22),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # ephemeral mcp_server × peers
    add_box(s, Inches(8.70), Inches(1.22), Inches(4.35), Inches(0.85),
            fill=RGBColor(0x1a,0x14,0x05), line=ORANGE)
    add_text(s, "mcp_server.py  × peers  (ephemeral)", Inches(8.80), Inches(1.30),
             Inches(4.15), Inches(0.30), size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, "subprocess.run()  ·  query_repo(feature_request)  ·  exits after call",
             Inches(8.80), Inches(1.62), Inches(4.15), Inches(0.22),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # score annotation
    add_box(s, Inches(8.70), Inches(2.18), Inches(4.35), Inches(0.65),
            fill=RGBColor(0x14,0x10,0x03), line=RGBColor(0x44,0x33,0x00))
    add_text(s, "score = min(len(response) / 800,  1.0)\n"
                "0.0 if '(no relevant knowledge found)'  ·  top-2 returned",
             Inches(8.80), Inches(2.22), Inches(4.15), Inches(0.55),
             size=8, color=ORANGE, align=PP_ALIGN.CENTER)

    # top-2 result annotation (return path label)
    add_box(s, Inches(3.20), Inches(2.18), Inches(4.30), Inches(0.65),
            fill=RGBColor(0x03,0x12,0x07), line=RGBColor(0x1a,0x40,0x1a))
    add_text(s, "③ returns: [repo-b, repo-c]  (top-2 by relevance score)\n"
                "test_mcp.py notes selected repos for Step 2",
             Inches(3.30), Inches(2.22), Inches(4.10), Inches(0.55),
             size=8, color=GREEN, align=PP_ALIGN.CENTER)

    # ── ROUTING connectors ─────────────────────────────────────────────────────
    # ① test_mcp → router
    conn(s, 2.38, 1.645, 3.20, 1.645, GREEN, ctype=2)
    seq_badge(s, 1, 2.79, 1.38, GREEN)
    add_text(s, "subprocess.run(router)  ·  tools/call get_relevant_repos",
             Inches(2.00), Inches(1.18), Inches(1.90), Inches(0.22),
             size=7, color=GREEN, align=PP_ALIGN.CENTER)

    # ② router → ephemeral mcp_server
    conn(s, 7.50, 1.645, 8.70, 1.645, ORANGE, ctype=2)
    seq_badge(s, 2, 8.10, 1.38, ORANGE)
    add_text(s, "subprocess.run(mcp_server) × R-1 peers",
             Inches(7.00), Inches(1.18), Inches(2.20), Inches(0.22),
             size=7, color=ORANGE, align=PP_ALIGN.CENTER)

    # ── QUERY boxes ────────────────────────────────────────────────────────────
    # test_mcp.py (query phase)
    add_box(s, Inches(0.28), Inches(3.78), Inches(2.10), Inches(0.85),
            fill=RGBColor(0x0d,0x1f,0x0d), line=GREEN)
    add_text(s, "test_mcp.py", Inches(0.38), Inches(3.86), Inches(1.90), Inches(0.36),
             size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, "_mcp_call() × 3 repos", Inches(0.38), Inches(4.22), Inches(1.90), Inches(0.22),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # mcp_server.py × 3
    add_box(s, Inches(3.20), Inches(3.78), Inches(3.10), Inches(0.85),
            fill=RGBColor(0x0d,0x1f,0x35), line=BLUE)
    add_text(s, "mcp_server.py  × 3", Inches(3.30), Inches(3.86), Inches(2.90), Inches(0.30),
             size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "rag.query(feature_request)  ·  ephemeral subprocess",
             Inches(3.30), Inches(4.18), Inches(2.90), Inches(0.22),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # Embedding Model (query embedding)
    add_box(s, Inches(7.50), Inches(3.78), Inches(3.30), Inches(0.85),
            fill=RGBColor(0x10,0x18,0x25), line=BLUE)
    add_text(s, "Embedding Model", Inches(7.60), Inches(3.86), Inches(3.10), Inches(0.30),
             size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "_embed([feature_request])  →  384-dim query vector",
             Inches(7.60), Inches(4.18), Inches(3.10), Inches(0.22),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # ChromaDB (similarity search)
    add_box(s, Inches(7.50), Inches(4.90), Inches(5.55), Inches(1.55),
            fill=RGBColor(0x0d,0x1f,0x12), line=GREEN)
    add_text(s, "ChromaDB  ·  .chroma_db/",
             Inches(7.60), Inches(4.97), Inches(5.35), Inches(0.30),
             size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, "docs_collection.query(query_embeddings, n_results=top_k)  →  top-k doc chunks + L2 distances\n"
                "code_collection.query(query_embeddings, n_results=top_k)  →  top-k code chunks + L2 distances\n"
                "keyword_search(question)  →  additional substring-matched chunks\n"
                "similarity_threshold filter  →  drop chunks with L2 dist > threshold",
             Inches(7.62), Inches(5.32), Inches(5.31), Inches(1.08),
             size=9, color=GREY)

    # RELEVANT KNOWLEDGE result box
    add_box(s, Inches(3.20), Inches(5.25), Inches(3.80), Inches(0.85),
            fill=RGBColor(0x1f,0x12,0x00), line=ORANGE)
    add_text(s, "RELEVANT KNOWLEDGE:", Inches(3.30), Inches(5.32), Inches(3.60), Inches(0.28),
             size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, "[From documentation] ... + [From source code] ...\nbuffered in proc.stdout → returned to test_mcp.py",
             Inches(3.30), Inches(5.62), Inches(3.60), Inches(0.42),
             size=8, color=GREY, align=PP_ALIGN.CENTER)

    # feature_analysis.md
    add_box(s, Inches(0.28), Inches(5.25), Inches(2.10), Inches(0.85),
            fill=RGBColor(0x1a,0x1a,0x0a), line=ORANGE)
    add_text(s, "feature_analysis.md", Inches(0.38), Inches(5.33), Inches(1.90), Inches(0.32),
             size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, "routing scores + all repo knowledge + instructions for Copilot",
             Inches(0.38), Inches(5.67), Inches(1.90), Inches(0.38),
             size=7, color=GREY, align=PP_ALIGN.CENTER)

    # ── QUERY connectors ───────────────────────────────────────────────────────
    # ④ test_mcp → mcp_server
    conn(s, 2.38, 4.205, 3.20, 4.205, BLUE, ctype=2)
    seq_badge(s, 4, 2.79, 3.95, BLUE)
    add_text(s, "subprocess.run(mcp_server)  ·  tools/call query_repo",
             Inches(2.00), Inches(3.74), Inches(1.90), Inches(0.22),
             size=7, color=BLUE, align=PP_ALIGN.CENTER)

    # ⑤ mcp_server → embedding
    conn(s, 6.30, 4.205, 7.50, 4.205, BLUE, ctype=2)
    seq_badge(s, 5, 6.90, 3.95, BLUE)
    add_text(s, "_embed([query])", Inches(6.20), Inches(3.74), Inches(1.90), Inches(0.22),
             size=7, color=BLUE, align=PP_ALIGN.CENTER)

    # ⑥ embedding → chromadb (straight down)
    conn(s, 9.15, 4.63, 9.15, 4.90, GREEN, ctype=1)
    seq_badge(s, 6, 9.38, 4.76, GREEN)
    add_text(s, "query_vector", Inches(9.55), Inches(4.70), Inches(1.40), Inches(0.20),
             size=7, color=GREEN)

    # ⑦ chromadb → RELEVANT KNOWLEDGE (elbow left)
    conn(s, 7.50, 5.675, 7.00, 5.675, ORANGE, ctype=2)
    seq_badge(s, 7, 7.15, 5.42, ORANGE)
    add_text(s, "top-k chunks + L2 dists", Inches(5.80), Inches(5.38), Inches(1.90), Inches(0.22),
             size=7, color=ORANGE, align=PP_ALIGN.CENTER)

    # RELEVANT KNOWLEDGE → test_mcp (horizontal left)
    conn(s, 3.20, 5.675, 2.38, 5.675, ORANGE, ctype=2)
    add_text(s, "RELEVANT KNOWLEDGE", Inches(2.00), Inches(5.38), Inches(1.90), Inches(0.22),
             size=7, color=ORANGE, align=PP_ALIGN.CENTER)

    # test_mcp → feature_analysis.md (straight down)
    conn(s, 1.33, 4.63, 1.33, 5.25, ORANGE, ctype=1)
    add_text(s, "writes", Inches(1.48), Inches(4.93), Inches(0.90), Inches(0.20),
             size=7, color=ORANGE)


def slide_vscode_flow(prs):
    """VS Code session-based orchestration: 3 repos + 1 orchestrator."""
    s = new_slide(prs)

    add_text(s, "DATA FLOW — VS CODE SESSION  ·  3 REPOS + ORCHESTRATOR",
             Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.26),
             size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "Permanent MCP Servers  ·  Streaming stdio Pipes  ·  Ephemeral Routing Subprocesses",
             Inches(0.4), Inches(0.35), Inches(12.5), Inches(0.42),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hline(s, Inches(0.83))

    # ── 4 VS Code window columns ───────────────────────────────────────────────
    win_names  = ["VS Code\nrepo-a", "VS Code\nrepo-b", "VS Code\nrepo-c", "VS Code\norchestrator"]
    win_colors = [BLUE, GREEN, PURPLE, ORANGE]
    win_fills  = [RGBColor(0x0d,0x1f,0x35), RGBColor(0x0d,0x1f,0x12),
                  RGBColor(0x2d,0x1f,0x47), RGBColor(0x22,0x14,0x00)]
    col_x = [0.22, 3.40, 6.58, 9.76]   # left edge of each column
    col_w = 2.95

    server_labels = [
        "router_mcp_server",
        "repo-a/mcp_server",
        "repo-b/mcp_server",
        "repo-c/mcp_server",
    ]
    srv_colors = [GREEN, BLUE, GREEN, PURPLE]

    for ci, (wx, wname, wcol, wfill) in enumerate(zip(col_x, win_names, win_colors, win_fills)):
        # VS Code window header box
        add_box(s, Inches(wx), Inches(0.92), Inches(col_w), Inches(0.56),
                fill=wfill, line=wcol)
        add_text(s, wname, Inches(wx+0.08), Inches(0.95), Inches(col_w-0.16), Inches(0.50),
                 size=11, bold=True, color=wcol, align=PP_ALIGN.CENTER)

        # 4 permanent MCP server process boxes per window
        for si, (slabel, scol) in enumerate(zip(server_labels, srv_colors)):
            sy = 1.65 + si * 0.85
            add_box(s, Inches(wx+0.08), Inches(sy), Inches(col_w-0.16), Inches(0.68),
                    fill=RGBColor(0x12,0x17,0x1f), line=scol)
            add_text(s, slabel, Inches(wx+0.16), Inches(sy+0.08), Inches(col_w-0.32), Inches(0.26),
                     size=9, bold=True, color=scol, align=PP_ALIGN.CENTER)
            add_text(s, "permanent · stdio pipe",
                     Inches(wx+0.16), Inches(sy+0.36), Inches(col_w-0.32), Inches(0.22),
                     size=7, color=GREY, align=PP_ALIGN.CENTER)

    # Process count annotation
    add_box(s, Inches(0.22), Inches(5.18), Inches(13.00), Inches(0.38),
            fill=RGBColor(0x10,0x12,0x18), line=BORDER)
    add_text(s, "4 windows  ×  4 servers each  =  16 permanent Python processes  "
                "(+ ~8 VS Code/Electron procs per window = ~48 total OS processes)",
             Inches(0.35), Inches(5.22), Inches(12.75), Inches(0.28),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Query flow section ─────────────────────────────────────────────────────
    add_box(s, Inches(0.22), Inches(5.65), Inches(13.00), Inches(1.65),
            fill=RGBColor(0x0a,0x12,0x0a), line=RGBColor(0x1a,0x30,0x1a))
    add_text(s, "When Copilot in repo-a window types a feature request:",
             Inches(0.35), Inches(5.68), Inches(12.75), Inches(0.24),
             size=9, bold=True, color=GREEN)

    # flow boxes in the bottom strip
    flow = [
        (0.35,  "Copilot\n(repo-a)", PURPLE, RGBColor(0x2d,0x1f,0x47)),
        (2.20,  "router_mcp_server\n(permanent proc)", GREEN, RGBColor(0x0d,0x1f,0x12)),
        (5.30,  "ephemeral\nmcp_server × peers", ORANGE, RGBColor(0x22,0x14,0x00)),
        (8.30,  "repo-a\nmcp_server", BLUE, RGBColor(0x0d,0x1f,0x35)),
        (10.40, "repo-b/c\nmcp_server", GREEN, RGBColor(0x0d,0x1f,0x12)),
    ]
    fh = 1.08
    for fx, flabel, fcol, ffill in flow:
        add_box(s, Inches(fx), Inches(5.98), Inches(1.70), Inches(fh),
                fill=ffill, line=fcol)
        add_text(s, flabel, Inches(fx+0.08), Inches(5.98+fh*0.18),
                 Inches(1.54), Inches(fh*0.65),
                 size=9, bold=True, color=fcol, align=PP_ALIGN.CENTER)

    # flow arrows in the bottom strip
    arrow_data = [
        (2.05, 6.52, 2.20, 6.52, GREEN,  "①\nstdio pipe",   1.98, 5.82),
        (4.10, 6.52, 5.30, 6.52, ORANGE, "② subprocess.run\n× peers (ephemeral)", 4.05, 5.82),
        (7.10, 6.52, 8.30, 6.52, BLUE,   "③ top-2 repos\nselected",               7.15, 5.82),
        (9.15, 6.52,10.40, 6.52, GREEN,  "④ query_repo\nstdio pipe",              9.10, 5.82),
    ]
    for ax1, ay1, ax2, ay2, acol, alb, ltx, lty in arrow_data:
        conn(s, ax1, ay1, ax2, ay2, acol, ctype=2)
        add_text(s, alb, Inches(ltx), Inches(lty), Inches(1.20), Inches(0.26),
                 size=6.5, color=acol, align=PP_ALIGN.CENTER)

    # ChromaDB access note
    add_text(s, "Each permanent mcp_server holds open its repo's .chroma_db/ (SQLite WAL).  "
                "Queries are concurrent reads — safe.  "
                "Router uses ephemeral subprocesses; does NOT reuse the permanent server processes.",
             Inches(0.35), Inches(7.12), Inches(12.75), Inches(0.30),
             size=8, color=GREY, align=PP_ALIGN.CENTER)


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

slide_title(prs)
slide_problem(prs)
slide_solution(prs)
slide_architecture(prs)
slide_tech_stack(prs)
slide_workflow(prs)
slide_rag(prs)
slide_mcp(prs)
slide_rag_init_flow(prs)
slide_rag_query_flow(prs)
slide_vscode_flow(prs)
slide_setup(prs)
slide_demo(prs)

out = Path(__file__).parent / "slides.pptx"
prs.save(str(out))
print(f"Saved: {out}  ({out.stat().st_size // 1024} KB,  {len(prs.slides)} slides)")

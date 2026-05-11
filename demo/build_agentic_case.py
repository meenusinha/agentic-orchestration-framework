"""
Generates demo/agentic_case.pptx — 5-slide executive presentation.
Why Agentic Orchestration beats a monolithic LLM: harness, evaluation, security.
Run from repo root:  python demo/build_agentic_case.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours (same palette as build_pptx.py) ───────────────────────────────────
BG        = RGBColor(0x0d, 0x11, 0x17)
BLUE      = RGBColor(0x58, 0xa6, 0xff)
WHITE     = RGBColor(0xf0, 0xf6, 0xfc)
GREY      = RGBColor(0x8b, 0x94, 0x9e)
GREEN     = RGBColor(0x3f, 0xb9, 0x50)
ORANGE    = RGBColor(0xf0, 0x88, 0x3e)
RED       = RGBColor(0xff, 0x7b, 0x72)
PURPLE    = RGBColor(0xbc, 0x8c, 0xff)
CARD_BG   = RGBColor(0x16, 0x1b, 0x22)
BORDER    = RGBColor(0x30, 0x36, 0x3d)
DARK_BLUE = RGBColor(0x1f, 0x3a, 0x5f)
MID_BLUE  = RGBColor(0x2d, 0x59, 0x86)
DARK_RED  = RGBColor(0x4a, 0x1a, 0x1a)
DARK_GREEN= RGBColor(0x1a, 0x3a, 0x1a)
DARK_PURP = RGBColor(0x2d, 0x1f, 0x4a)

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def txt(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return txb

def box(slide, x, y, w, h, fill=CARD_BG, line=BORDER, lw=Emu(9525)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = lw
    return s

def boxt(slide, text, x, y, w, h, size=12, bold=False, color=WHITE,
         fill=CARD_BG, line=BORDER, align=PP_ALIGN.LEFT):
    box(slide, x, y, w, h, fill=fill, line=line)
    pad = Inches(0.13)
    txt(slide, text, x+pad, y+pad, w-pad*2, h-pad*2,
        size=size, bold=bold, color=color, align=align)

def hline(slide, y, x1=0.5, x2=12.83):
    ln = slide.shapes.add_connector(1, Inches(x1), y, Inches(x2), y)
    ln.line.color.rgb = BORDER
    ln.line.width = Emu(9525)

def tag(slide, text, x, y):
    """Pill-style category tag."""
    txt(slide, text, x, y, Inches(3.5), Inches(0.28),
        size=9, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

def icon_bullet(slide, icon, text, x, y, w=Inches(5.5), icon_color=BLUE, size=13):
    txt(slide, icon, x, y, Inches(0.4), Inches(0.35),
        size=size, bold=True, color=icon_color)
    txt(slide, text, x+Inches(0.38), y, w-Inches(0.38), Inches(0.4),
        size=size, color=WHITE)

def slide_number(slide, n, total=5):
    txt(slide, f"{n} / {total}", Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
        size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — The Case: Parallel Teams, Bottleneck AI
# ─────────────────────────────────────────────────────────────────────────────
def slide_the_case(prs):
    s = new_slide(prs)

    tag(s, "THE PROBLEM", Inches(0.5), Inches(0.28))
    txt(s, "Your Teams Work in Parallel.\nYour AI Doesn't.",
        Inches(0.5), Inches(0.55), Inches(12.3), Inches(1.5),
        size=36, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    hline(s, Inches(1.95))

    # Left column — problem framing
    txt(s, "The Reality",
        Inches(0.5), Inches(2.1), Inches(4.0), Inches(0.35),
        size=13, bold=True, color=WHITE)
    reality = [
        ("◆", "5+ sub-function teams commit simultaneously"),
        ("◆", "Each team needs AI help right now, in parallel"),
        ("◆", "Each domain has its own code, physics, constraints"),
        ("◆", "Any change can invalidate an AI trained on the whole stack"),
    ]
    for i, (ic, t) in enumerate(reality):
        icon_bullet(s, ic, t, Inches(0.5), Inches(2.55 + i*0.48),
                    w=Inches(4.2), icon_color=ORANGE)

    # Centre divider with label
    for yy in [2.1, 2.55, 3.03, 3.51, 3.99]:
        txt(s, "→", Inches(5.0), Inches(yy), Inches(0.5), Inches(0.4),
            size=16, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    # Right column — 4 failure modes as cards
    txt(s, "4 Critical Failure Modes of a Monolithic LLM",
        Inches(5.5), Inches(2.1), Inches(7.3), Inches(0.35),
        size=13, bold=True, color=RED)

    failures = [
        (RED,    "① Retraining Bottleneck",
                 "One team's commit → 3–4 week retrain → all others blocked"),
        (ORANGE, "② Knowledge Conflicts",
                 "Optics ≠ Control Theory ≠ Sensor Physics — a single model degrades at all three"),
        (PURPLE, "③ No Auditability",
                 "Mixed-domain training = can't trace a suggestion to its source — fatal for safety-critical code"),
        (GREY,   "④ Context Saturation",
                 "All designs + specs + code > any LLM context — model drowns in irrelevant noise"),
    ]
    for i, (col, title, body) in enumerate(failures):
        bx = Inches(5.5)
        by = Inches(2.55 + i*0.92)
        box(s, bx, by, Inches(7.3), Inches(0.84),
            fill=RGBColor(0x1a, 0x1a, 0x22), line=col, lw=Emu(19050))
        txt(s, title, bx+Inches(0.13), by+Inches(0.08),
            Inches(7.0), Inches(0.3), size=11, bold=True, color=col)
        txt(s, body, bx+Inches(0.13), by+Inches(0.38),
            Inches(7.0), Inches(0.38), size=10, color=GREY)

    # Bottom call-out
    hline(s, Inches(6.68))
    txt(s, "Bottom line:  Decomposing your code solved the parallel-team problem.  "
           "Decomposing your AI solves it again — at the intelligence layer.",
        Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number(s, 1)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Solution: Distributed Agents vs Monolithic LLM
# ─────────────────────────────────────────────────────────────────────────────
def slide_solution(prs):
    s = new_slide(prs)

    tag(s, "THE SOLUTION", Inches(0.5), Inches(0.28))
    txt(s, "One Agent Per Sub-function, All Coordinated",
        Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
        size=30, bold=True, color=BLUE)
    hline(s, Inches(1.2))

    # ── Left: Monolithic (crossed out feel) ──────────────────────────────────
    box(s, Inches(0.4), Inches(1.35), Inches(4.1), Inches(4.85),
        fill=RGBColor(0x20, 0x10, 0x10), line=RED, lw=Emu(19050))
    txt(s, "✗  Monolithic LLM", Inches(0.55), Inches(1.45),
        Inches(3.8), Inches(0.4), size=13, bold=True, color=RED)
    mono_pts = [
        "Trained on entire codebase",
        "One stale update blocks everyone",
        "Cannot isolate domain expertise",
        "Suggestions untraceable by domain",
        "Context window overwhelmed",
        "3–4 week retrain per change",
    ]
    for i, pt in enumerate(mono_pts):
        txt(s, f"  ✗  {pt}", Inches(0.6), Inches(1.98 + i*0.55),
            Inches(3.75), Inches(0.48), size=11, color=RGBColor(0xcc, 0x60, 0x60))

    # ── Middle arrow ──────────────────────────────────────────────────────────
    txt(s, "vs", Inches(4.6), Inches(3.6), Inches(0.6), Inches(0.5),
        size=20, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    # ── Right: Agentic architecture diagram ──────────────────────────────────
    box(s, Inches(5.3), Inches(1.35), Inches(7.5), Inches(4.85),
        fill=RGBColor(0x10, 0x20, 0x10), line=GREEN, lw=Emu(19050))
    txt(s, "✓  Distributed Agentic Orchestration", Inches(5.45), Inches(1.45),
        Inches(7.2), Inches(0.4), size=13, bold=True, color=GREEN)

    # Orchestrator box
    boxt(s, "Orchestrator\n(cross-domain routing)",
         Inches(7.8), Inches(2.0), Inches(2.5), Inches(0.72),
         size=10, bold=True, color=WHITE, fill=DARK_BLUE, line=BLUE)

    # Agent boxes
    agents = [
        ("Wavefront\nAgent", 5.5, 3.15),
        ("Stage Control\nAgent", 7.9, 3.15),
        ("Metrology\nAgent", 10.3, 3.15),
    ]
    for name, ax, ay in agents:
        boxt(s, name, Inches(ax), Inches(ay), Inches(2.0), Inches(0.72),
             size=10, bold=True, color=WHITE, fill=DARK_GREEN, line=GREEN)

    # RAG tag under each agent
    for ax, ay in [(5.5, 3.87), (7.9, 3.87), (10.3, 3.87)]:
        boxt(s, "Domain LLM\n+ RAG Index",
             Inches(ax), Inches(ay), Inches(2.0), Inches(0.58),
             size=8, color=GREY, fill=RGBColor(0x0f,0x18,0x0f), line=BORDER)

    # Connectors: orchestrator → agents
    for ax in [6.5, 9.0, 11.3]:
        c = s.shapes.add_connector(1, Inches(9.05), Inches(2.72), Inches(ax), Inches(3.15))
        c.line.color.rgb = GREEN
        c.line.width = Emu(12700)

    # Bullet points for agentic
    agentic_pts = [
        "Deep narrow expertise per sub-function",
        "RAG index updates on each commit (minutes)",
        "Agents run in parallel — zero blocking",
        "Every suggestion traceable to its knowledge source",
        "Orchestrator handles cross-domain features",
        "Scales: add a new sub-function → add one agent",
    ]
    for i, pt in enumerate(agentic_pts):
        txt(s, f"  ✓  {pt}", Inches(5.45), Inches(4.6 + i*0.47),
            Inches(7.2), Inches(0.42), size=10, color=RGBColor(0x60, 0xcc, 0x80))

    # Research backing bar
    hline(s, Inches(6.68))
    txt(s, "Research backing:  MetaGPT (Hong et al. 2023) · Self-Organized Agents (Ishibashi & Nishimura 2024) · "
           "Cross-Team Orchestration (Du et al. ACL 2025) · Anthropic: Building Effective Agents (2024)",
        Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55),
        size=9, color=GREY, align=PP_ALIGN.CENTER)

    slide_number(s, 2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Agent Harness
# ─────────────────────────────────────────────────────────────────────────────
def slide_harness(prs):
    s = new_slide(prs)

    tag(s, "GOVERNANCE", Inches(0.5), Inches(0.28))
    txt(s, "The Agent Harness: Safety & Quality by Design",
        Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
        size=30, bold=True, color=BLUE)
    hline(s, Inches(1.2))

    # Definition card
    box(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.88),
        fill=DARK_BLUE, line=BLUE, lw=Emu(19050))
    txt(s, "What is a harness?",
        Inches(0.6), Inches(1.4), Inches(3.5), Inches(0.35),
        size=12, bold=True, color=BLUE)
    txt(s, "A controlled runtime wrapper around each agent that enforces tool permissions, "
           "runs standard checks, records full traces, and produces one decision report — "
           "before any output reaches a human or downstream system.",
        Inches(0.6), Inches(1.73), Inches(12.1), Inches(0.42),
        size=11, color=WHITE)

    # 4 harness responsibilities
    cols = [
        (BLUE,   "Tool\nPermissions",
                 "Each agent can only call\nthe MCP servers it owns.\nNo cross-boundary reads\nwithout explicit policy."),
        (GREEN,  "Standard\nChecks",
                 "Physics consistency,\nregression tests, linting,\ntype checks — all run\nautomatically per change."),
        (ORANGE, "Trace\nLogging",
                 "Every tool call, RAG query,\nand model response logged.\nFull audit trail from request\nto recommendation."),
        (PURPLE, "Decision\nReport",
                 "One scorecard per change:\npassed checks, risk flags,\ntool-call trace, source-\nartifact citation."),
    ]
    for i, (col, title, body) in enumerate(cols):
        cx = Inches(0.4 + i * 3.15)
        cy = Inches(2.42)
        box(s, cx, cy, Inches(3.0), Inches(2.5),
            fill=RGBColor(0x12, 0x17, 0x1f), line=col, lw=Emu(19050))
        txt(s, title, cx+Inches(0.12), cy+Inches(0.12),
            Inches(2.76), Inches(0.62), size=13, bold=True, color=col,
            align=PP_ALIGN.CENTER)
        hline(s, cy+Inches(0.78))
        txt(s, body, cx+Inches(0.12), cy+Inches(0.85),
            Inches(2.76), Inches(1.52), size=11, color=GREY,
            align=PP_ALIGN.LEFT)

    # Implementation strip
    hline(s, Inches(5.1))
    txt(s, "Implementation stack",
        Inches(0.5), Inches(5.18), Inches(3.0), Inches(0.32),
        size=11, bold=True, color=WHITE)

    stack = [
        (BLUE,   "LangGraph /\nLangChain",  "Stateful agent workflow,\nconditional routing, retries"),
        (GREEN,  "MCP Servers",             "Scoped tool access per\nrepo — stdio JSON-RPC"),
        (ORANGE, "CI/CD Pipeline",          "Harness runs in PR gate;\nblocks merge on failure"),
        (PURPLE, "Governance\nControls",    "Access control, encryption,\naudit log retention"),
    ]
    for i, (col, title, body) in enumerate(stack):
        sx = Inches(0.4 + i * 3.15)
        sy = Inches(5.5)
        boxt(s, f"{title}\n{body}", sx, sy, Inches(3.0), Inches(0.9),
             size=10, color=WHITE, fill=RGBColor(0x10,0x15,0x1c), line=col)

    hline(s, Inches(6.68))
    txt(s, "\"No unharnessed output ships.\"  — AI accelerates; the harness governs; humans decide.",
        Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.5),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number(s, 3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Evaluation Framework
# ─────────────────────────────────────────────────────────────────────────────
def slide_evaluation(prs):
    s = new_slide(prs)

    tag(s, "EVALUATION", Inches(0.5), Inches(0.28))
    txt(s, "Two Gates Before Any Output Reaches Production",
        Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
        size=30, bold=True, color=BLUE)
    hline(s, Inches(1.2))

    # Gate 1 — Automated
    box(s, Inches(0.4), Inches(1.35), Inches(5.9), Inches(4.7),
        fill=RGBColor(0x10, 0x1a, 0x10), line=GREEN, lw=Emu(22225))
    txt(s, "Gate 1 — Automated Evaluation", Inches(0.6), Inches(1.45),
        Inches(5.5), Inches(0.4), size=14, bold=True, color=GREEN)
    txt(s, "Runs on every commit. Blocks Gate 2 if it fails.",
        Inches(0.6), Inches(1.85), Inches(5.5), Inches(0.32),
        size=10, color=GREY)
    hline(s, Inches(2.17), x1=0.5, x2=6.2)

    auto_pts = [
        (GREEN,  "HELM metrics",      "Accuracy · Robustness · Consistency\n(Bommasani et al. 2023 — industry standard)"),
        (BLUE,   "Domain checks",     "Physics consistency, performance benchmarks,\nregression vs. reference artifacts"),
        (ORANGE, "Risk flags",        "Hallucination detector, confidence score,\nout-of-domain query detection"),
        (PURPLE, "Scorecard output",  "Pass/Fail threshold, tool-call trace,\nsource-artifact citation per suggestion"),
    ]
    for i, (col, title, body) in enumerate(auto_pts):
        px = Inches(0.6)
        py = Inches(2.27 + i * 0.85)
        txt(s, f"▸  {title}", px, py, Inches(5.5), Inches(0.3),
            size=11, bold=True, color=col)
        txt(s, body, px+Inches(0.25), py+Inches(0.3), Inches(5.2), Inches(0.45),
            size=10, color=GREY)

    # Gate 2 — Human
    box(s, Inches(6.5), Inches(1.35), Inches(6.4), Inches(4.7),
        fill=RGBColor(0x1a, 0x14, 0x10), line=ORANGE, lw=Emu(22225))
    txt(s, "Gate 2 — Human-in-Loop Review", Inches(6.7), Inches(1.45),
        Inches(6.0), Inches(0.4), size=14, bold=True, color=ORANGE)
    txt(s, "Final decision. Engineer sees WHY, not just WHAT.",
        Inches(6.7), Inches(1.85), Inches(6.0), Inches(0.32),
        size=10, color=GREY)
    hline(s, Inches(2.17), x1=6.6, x2=12.8)

    human_pts = [
        (ORANGE, "Evaluated results review",   "Engineer reviews scorecard + source trace\nbefore accepting any AI recommendation"),
        (RED,    "Design sign-off",             "Architecture and design decisions always\nrequire explicit engineer approval"),
        (PURPLE, "Cross-boundary features",     "Any feature touching >1 sub-function\nescalates to system lead review"),
        (BLUE,   "Compliance & IP checks",      "Security, IP exposure, licence compliance\nreviewed before merge"),
    ]
    for i, (col, title, body) in enumerate(human_pts):
        px = Inches(6.7)
        py = Inches(2.27 + i * 0.85)
        txt(s, f"▸  {title}", px, py, Inches(6.0), Inches(0.3),
            size=11, bold=True, color=col)
        txt(s, body, px+Inches(0.25), py+Inches(0.3), Inches(5.75), Inches(0.45),
            size=10, color=GREY)

    # Arrow between gates
    arr = s.shapes.add_connector(1, Inches(6.32), Inches(3.7), Inches(6.5), Inches(3.7))
    arr.line.color.rgb = WHITE
    arr.line.width = Emu(19050)

    # Bottom bar
    hline(s, Inches(6.22))
    metric_cols = [
        (GREEN,  "Automated",  "Accuracy,\nRobustness,\nConsistency"),
        (BLUE,   "Domain",     "Physics,\nPerformance,\nRegression"),
        (ORANGE, "Risk",       "Hallucination,\nConfidence,\nOOD detection"),
        (PURPLE, "Trace",      "Full tool-call\nand source-\nartifact log"),
        (RED,    "Human gate", "Sign-off on\ndesign decisions\n+ compliance"),
    ]
    for i, (col, title, body) in enumerate(metric_cols):
        mx = Inches(0.4 + i * 2.5)
        my = Inches(6.3)
        box(s, mx, my, Inches(2.35), Inches(0.95),
            fill=RGBColor(0x12,0x14,0x1a), line=col, lw=Emu(12700))
        txt(s, title, mx+Inches(0.1), my+Inches(0.05),
            Inches(2.1), Inches(0.28), size=10, bold=True, color=col)
        txt(s, body, mx+Inches(0.1), my+Inches(0.32),
            Inches(2.1), Inches(0.55), size=8, color=GREY)

    slide_number(s, 4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Security: 4-Layer IP Protection
# ─────────────────────────────────────────────────────────────────────────────
def slide_security(prs):
    s = new_slide(prs)

    tag(s, "SECURITY & IP", Inches(0.5), Inches(0.28))
    txt(s, "Four-Layer IP Protection: Risk Distributed, Not Concentrated",
        Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
        size=28, bold=True, color=BLUE)
    hline(s, Inches(1.2))

    layers = [
        # (border_color, fill_color, layer_num, title, col1_title, col1_body, col2_title, col2_body)
        (
            BLUE, RGBColor(0x10,0x18,0x28),
            "Layer 1",
            "Specialist Models Per Sub-Function AND Per Task",
            "Code Model  (Codex-class)",
            "Knows HOW to write domain-specific code.\nNot WHY it matters. Not cross-function APIs.\n→ Codex: 28.8% HumanEval vs GPT-3: 0%\n   (Chen et al. 2021)",
            "Reasoning Model  (Chain-of-thought)",
            "Knows WHY: physics constraints, domain logic.\nCannot generate code syntax alone.\n→ Compromise ≠ code generation capability\n   (OpenAI o1, Anthropic Extended Reasoning)",
        ),
        (
            GREEN, RGBColor(0x10,0x22,0x14),
            "Layer 2",
            "Model Protection = Production Code Protection",
            "Technical controls",
            "Versioning + rollback · Access control\nEncryption at rest & in transit\nAudit logging of every prediction",
            "Deployment options",
            "Air-gapped on-prem: models never leave\nyour infrastructure · No cloud API calls\nfor sensitive sub-function code",
        ),
        (
            ORANGE, RGBColor(0x22,0x18,0x10),
            "Layer 3",
            "RAG for Live Knowledge — Stateless, Portable, Updateable",
            "Models are stateless",
            "Code + specs live in YOUR vector DB.\nUpdate the database → all agents see new\nknowledge instantly. No retraining needed.",
            "Security benefit",
            "Stolen model = code skeleton generator\nwithout your proprietary knowledge.\nKnowledge and model are separated.",
        ),
        (
            PURPLE, RGBColor(0x1c,0x14,0x2a),
            "Layer 4",
            "Scoped MCP Permissions + Human Final Gate",
            "Scoped access",
            "Each model accesses only its sub-function's\nMCP servers. Cross-boundary decisions always\nrequire explicit human approval.",
            "Attack surface",
            "Attacker must compromise: Code Model +\nReasoning Model + Orchestrator + RAG DB\n= far harder than one monolithic target.",
        ),
    ]

    for i, (bcol, fcol, lnum, title, c1t, c1b, c2t, c2b) in enumerate(layers):
        row_y = Inches(1.35 + i * 1.47)
        bx = Inches(0.4)
        row_h = Inches(1.38)

        box(s, bx, row_y, Inches(12.5), row_h, fill=fcol, line=bcol, lw=Emu(15876))

        # Layer badge
        boxt(s, lnum, bx+Inches(0.08), row_y+Inches(0.08),
             Inches(0.82), Inches(0.32), size=9, bold=True,
             color=BG, fill=bcol, line=bcol, align=PP_ALIGN.CENTER)

        # Title
        txt(s, title, bx+Inches(1.0), row_y+Inches(0.08),
            Inches(11.3), Inches(0.35), size=12, bold=True, color=bcol)

        # Col 1
        txt(s, c1t, bx+Inches(1.0), row_y+Inches(0.5),
            Inches(5.5), Inches(0.28), size=10, bold=True, color=WHITE)
        txt(s, c1b, bx+Inches(1.0), row_y+Inches(0.78),
            Inches(5.5), Inches(0.52), size=9, color=GREY)

        # Divider
        dv = s.shapes.add_connector(1, Inches(6.6), row_y+Inches(0.45),
                                     Inches(6.6), row_y+row_h-Inches(0.1))
        dv.line.color.rgb = BORDER
        dv.line.width = Emu(9525)

        # Col 2
        txt(s, c2t, bx+Inches(6.35), row_y+Inches(0.5),
            Inches(5.9), Inches(0.28), size=10, bold=True, color=WHITE)
        txt(s, c2b, bx+Inches(6.35), row_y+Inches(0.78),
            Inches(5.9), Inches(0.52), size=9, color=GREY)

    # Bottom
    hline(s, Inches(7.1))
    txt(s, "Outcome:  Risk is distributed across sub-functions, task types, and protection layers.  "
           "No single model is a complete mirror of your IP.",
        Inches(0.4), Inches(7.16), Inches(12.5), Inches(0.28),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number(s, 5)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_the_case(prs)
    slide_solution(prs)
    slide_harness(prs)
    slide_evaluation(prs)
    slide_security(prs)

    out = Path(__file__).parent / "agentic_case.pptx"
    prs.save(str(out))
    print(f"Saved {out}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
